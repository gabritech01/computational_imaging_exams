"""Builds the exam presentation (.pptx), following the structure required by the
course's presentation trace: title, description, methodology, implementation,
experiments, numerical results, conclusions, bibliography."""
import csv

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RED = RGBColor(0xB2, 0x2E, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_slide(prs, title, subtitle):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RED
    bg.line.fill.background()
    bg.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.5), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(title.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(subtitle):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
    return slide


def add_section_title(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER
    return tb


def add_content_slide(prs, title):
    slide = blank_slide(prs)
    add_section_title(slide, title)
    return slide


def add_bullets(slide, bullets, top=1.3, left=0.8, width=11.7, height=5.6, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if level == 0 else "-  ") + text
        p.level = level
        p.font.size = Pt(size - level * 2)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def add_table(slide, headers, rows, top=1.4, left=0.8, width=11.7, height=3.5):
    n_rows, n_cols = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
    return table


def read_comparison_csv():
    with open("results/tables/comparison_summary.csv") as f:
        return list(csv.DictReader(f))


def build():
    prs = new_pres()

    # 1. Title
    add_title_slide(
        prs,
        "Deblur & Denoise on KaoKore\nas an Inverse Problem",
        ["Computational Imaging 2025-26", "Gabriele Centonze"],
    )

    # 2. Description of the project
    s = add_content_slide(prs, "Description of the project")
    add_bullets(s, [
        ("Task: deblur + denoise, formulated as an inverse problem y = Ax + n", 0),
        ("A: Gaussian blur (sigma=2, kernel 9x9); n: additive Gaussian noise, 4 levels (0.005, 0.01, 0.05, 0.1)", 1),
        ("Dataset: KaoKore v1.3, 256x256 RGB portraits (9257 images)", 0),
        ("Official train / dev / test split from the dataset's own labels.csv: 7405 / 926 / 926", 1),
        ("Goal: implement and critically compare methods from different methodological families", 0),
        ("under the same experimental conditions (same degraded inputs for every method)", 1),
        ("3 of the 4 proposed methods (allowed for single-student projects):", 0),
        ("Variational (FISTA + Wavelet), Hybrid (PD-Net + TV), End-to-end (UNet)", 1),
    ])

    # 3. Methodology overview
    s = add_content_slide(prs, "Methodology — overview")
    add_bullets(s, [
        ("The inverse problem is ill-posed: A attenuates high spatial frequencies", 0),
        ("naive inversion amplifies noise uncontrollably -> regularization / a prior is required", 1),
        ("Three different ways of introducing that prior:", 0),
        ("Variational: an explicit hand-crafted prior (wavelet sparsity), no learning", 1),
        ("Hybrid: known physics (A, image gradient) + a few learned correction blocks", 1),
        ("End-to-end: the whole mapping y -> x is learned from data", 1),
        ("All three share the exact same forward operator and read the exact same", 0),
        ("pre-generated degraded PNG files for evaluation (fair comparison)", 1),
    ])

    # 4. Methodology: FISTA
    s = add_content_slide(prs, "Methodology — Variational: FISTA + Wavelet")
    add_image(s, "results/figures/formula_fista.png", 1.9, 1.3, height=0.75)
    add_bullets(s, [
        ("W: orthogonal wavelet transform", 0),
        ("Non-smooth but convex -> proximal gradient method (Beck & Teboulle, 2009)", 0),
        ("For orthogonal W, prox of the L1 term = soft-thresholding of wavelet coefficients", 0),
        ("Nesterov-style extrapolation: convergence rate O(1/k^2) instead of O(1/k)", 0),
        ("Step size 1/L: L = ||A^T A|| = 1 exactly, thanks to circular boundary conditions", 0),
        ("(A becomes a symmetric circulant operator, diagonalized by the Fourier transform)", 1),
    ], top=2.5)

    # 5. Methodology: PD-Net
    s = add_content_slide(prs, "Methodology — Hybrid: PD-Net + TV")
    add_image(s, "results/figures/formula_pdnet.png", 2.4, 1.3, height=0.75)
    add_bullets(s, [
        ("Unrolled Chambolle-Pock primal-dual algorithm", 0),
        ("Exact proximal / projection steps replaced by small learned CNN blocks", 0),
        ("Dual variable lives in gradient space (not data space): ties the network", 0),
        ("explicitly to the TV structure, unlike a generic learned primal-dual", 1),
        ("Data-fidelity gradient A^T(Ax-y) stays exact (not learned) at every iteration", 0),
        ("-> hybrid: known physics + learned correction, not a black box", 1),
    ], top=2.5)

    # 6. Methodology: UNet
    s = add_content_slide(prs, "Methodology — End-to-end: UNet")
    add_bullets(s, [
        ("Fully supervised: learns the direct mapping y -> x from training pairs", 0),
        ("Chosen over ViT / NAF-Net: the degradation is spatially local (small kernel,", 0),
        ("per-pixel noise) -> no need for long-range attention; simpler to implement well", 1),
        ("in the available time", 1),
        ("Encoder-decoder with skip connections: preserve fine spatial detail across", 0),
        ("the bottleneck, needed since input and output share almost all structure", 1),
        ("Residual learning (the network only learns the correction):", 0),
    ])
    add_image(s, "results/figures/formula_unet_residual.png", 2.9, 4.9, height=0.65)

    # 7. Implementation: FISTA
    s = add_content_slide(prs, "Implementation — FISTA + Wavelet")
    add_bullets(s, [
        ("Wavelet: Daubechies db4, 3 decomposition levels, mode='periodization'", 0),
        ("(periodization keeps the DWT exactly orthogonal, required for the prox to be exact)", 1),
        ("100 iterations per image, no learned parameters", 0),
        ("lambda tuned by grid search on the dev set (20 images), separately per noise level", 0),
        ("(0.005 -> 0.0025, 0.01 -> 0.0035, 0.05 -> 0.0175, 0.1 -> 0.07: grows with noise, as", 1),
        ("expected from the discrepancy principle)", 1),
    ])

    # 8. Implementation: PD-Net
    s = add_content_slide(prs, "Implementation — PD-Net")
    add_image(s, "results/figures/pdnet_diagram.png", 1.3, 1.3, width=10.7)
    add_bullets(s, [
        ("8 unrolled iterations, independent weights per iteration; CNN blocks: Conv3x3-LeakyReLU-Conv3x3", 0),
        ("4 specialized models (one per noise level), 700 training images, 12 epochs, Adam (lr=2e-4), L1 loss", 0),
    ], top=5.9, size=15)

    # 9. Implementation: UNet
    s = add_content_slide(prs, "Implementation — UNet")
    add_image(s, "results/figures/unet_diagram.png", 1.9, 1.3, width=9.6)
    add_bullets(s, [
        ("4 downsampling levels, base channels 48, GroupNorm (stable with small batch size)", 0),
        ("4 specialized models (one per noise level), 1000 training images, 20 epochs, Adam (lr=2e-4), L1 loss", 0),
    ], top=5.9, size=15)

    # 10. Introduce the experiments
    s = add_content_slide(prs, "Introduce the experiments")
    add_bullets(s, [
        ("Degradation: Gaussian blur sigma=2, kernel 9x9, circular boundary conditions", 0),
        ("Noise: additive Gaussian, std = noise level directly (0.005 / 0.01 / 0.05 / 0.1), on [0,1] images", 0),
        ("Degraded observations generated once, saved as 8-bit PNGs -> identical inputs for all methods", 0),
        ("Time-constrained scope decisions (declared explicitly, not hidden):", 0),
        ("evaluation on a fixed subset (80 test + 20 dev images) instead of the full 926/926", 1),
        ("3 of 4 methods (allowed for single-student projects); generative method (Diffusion+DPS) dropped", 1),
        ("4 specialized models per learned method, for a fair comparison against FISTA's per-level lambda", 1),
    ])

    # 11. Results table
    s = add_content_slide(prs, "Numerical results — summary table")
    rows_data = read_comparison_csv()
    table_rows = [[r["method"], r["noise_level"], f"{float(r['psnr_mean']):.2f}",
                   f"{float(r['ssim_mean']):.3f}", f"{float(r['avg_time_sec']):.4f}"]
                  for r in rows_data]
    add_table(s, ["Method", "Noise level", "PSNR (dB)", "SSIM", "Time/image (s)"], table_rows,
              top=1.3, height=5.7)

    # 12. Results plot
    s = add_content_slide(prs, "Numerical results — comparison plot")
    add_image(s, "results/figures/comparison_plot.png", 1.0, 1.5, width=11.3)

    # 13. Results visual comparison
    s = add_content_slide(prs, "Numerical results — visual comparison")
    add_image(s, "results/figures/composite_comparison.png", 3.1, 1.2, height=6.15)

    # 14. Conclusions
    s = add_content_slide(prs, "Conclusions")
    add_bullets(s, [
        ("UNet and PD-Net outperform FISTA by 2-3 dB PSNR at every noise level:", 0),
        ("a prior learned on KaoKore's own statistics beats a generic wavelet-sparsity prior", 1),
        ("UNet and PD-Net achieve near-identical accuracy, but PD-Net uses ~100x fewer parameters", 0),
        ("(302 KB vs 33 MB checkpoints): embedding known physics reduces what must be learned", 1),
        ("FISTA is ~40x slower per image (iterative optimization vs a single forward pass),", 0),
        ("but needs no training and is fully interpretable, unlike the two learned methods", 1),
        ("Limits: reduced evaluation subset and training set/epochs, due to time constraints", 0),
        ("Future work: full test set and training set, more epochs, add the 4th method", 0),
        ("(Diffusion + DPS), try NAF-Net as an alternative end-to-end architecture", 1),
    ])

    prs.save("presentation/presentation.pptx")
    print("saved presentation/presentation.pptx")


if __name__ == "__main__":
    build()
