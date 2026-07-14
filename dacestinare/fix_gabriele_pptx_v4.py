"""Fourth pass: swaps the embedded comparison plot for the regenerated one that
includes the degraded-observation baseline curve (the on-disk PNG was updated but the
pptx still had the old image embedded). Same position and size as the original."""
from pptx import Presentation
from pptx.util import Emu

PATH = "presentation/gabriele_centonze.pptx"
PLOT = "results/figures/comparison_plot.png"


def title_of(slide):
    best = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.top or 0
            if best is None or t < best[0]:
                best = (t, sh.text_frame.text)
    return best[1] if best else ""


def main():
    prs = Presentation(PATH)
    for slide in prs.slides:
        if title_of(slide).startswith("Confronto Prestazionale"):
            plot = None
            for sh in slide.shapes:
                if sh.shape_id == 224:
                    plot = sh
            left, top, width, height = plot.left, plot.top, plot.width, plot.height
            plot._element.getparent().remove(plot._element)
            slide.shapes.add_picture(PLOT, left, top, width, height)
            break
    prs.save(PATH)
    print("saved", PATH)


if __name__ == "__main__":
    main()
