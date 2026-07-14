"""Applies verified corrections to presentation/gabriele_centonze.pptx.

Every edit was checked against the actual shape positions first (see the
measurement pass in the session): appends only go into the LAST bullet shape of
a slide (nothing below it to collide with) or into fresh space below all
existing content; in-place replacements on shapes with siblings right below
them are kept short so autofit growth can't push into the next shape.
"""
import copy
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

PATH = "presentation/gabriele_centonze.pptx"
BACKUP = "presentation/old/gabriele_centonze_before_fix.pptx"

TITLE_FONT, TITLE_SIZE, TITLE_COLOR = "Merriweather", 24, RGBColor(0x10, 0xB9, 0x81)
BODY_FONT, BODY_SIZE, BODY_COLOR = "Open Sans", 14, RGBColor(0xCB, 0xD5, 0xE1)


def get_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"shape {shape_id} not found")


def set_run_text(slide, shape_id, para_idx, run_idx, text):
    get_shape(slide, shape_id).text_frame.paragraphs[para_idx].runs[run_idx].text = text


def append_run_text(slide, shape_id, para_idx, run_idx, extra):
    r = get_shape(slide, shape_id).text_frame.paragraphs[para_idx].runs[run_idx]
    r.text = r.text + extra


def style_run(run, size=BODY_SIZE, color=BODY_COLOR, font=BODY_FONT, italic=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    if italic is not None:
        run.font.italic = italic


def add_textbox(slide, left, top, width, height, paragraphs, size=BODY_SIZE):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        style_run(r, size=size)
        p.space_after = Pt(8)
    return tb


def duplicate_slide(prs, source_slide):
    """Copies a slide's shapes (with formatting) to a new slide at the end,
    remapping image relationship ids so pictures/backgrounds aren't broken."""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    rid_map = {}
    for rId, rel in source_slide.part.rels.items():
        if rel.reltype.endswith("slideLayout") or rel.reltype.endswith("notesSlide"):
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rId] = new_rid

    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        for el in new_el.iter():
            for attr in ("embed", "link"):
                val = el.get(qn("r:" + attr))
                if val in rid_map:
                    el.set(qn("r:" + attr), rid_map[val])
        new_slide.shapes._spTree.append(new_el)
    return new_slide


def set_paragraph_text(shape, para_idx, text, size=None, bold=None, color=None, font=None):
    p = shape.text_frame.paragraphs[para_idx]
    if not p.runs:
        p.add_run()
    r = p.runs[0]
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    r.text = text
    if size:
        r.font.size = Pt(size)
    if bold is not None:
        r.font.bold = bold
    if color:
        r.font.color.rgb = color
    if font:
        r.font.name = font


def main():
    shutil.copy(PATH, BACKUP)
    prs = Presentation(PATH)
    slides = list(prs.slides)

    # === Slide 2: Descrizione del Progetto ===
    # shape 96 (Metodi implementati) is the LAST bullet, bottom=4.99in, 2.5in free below.
    s2 = slides[1]
    add_textbox(s2, 0.75, 5.3, 11.6, 1.8, [
        "Nota sui numeri: split ufficiale citato sopra, ma per vincoli di tempo tuning "
        "e valutazione usano un sottoinsieme fisso (20 dev, 80 test); il training usa "
        "700-1000 immagini invece delle 7405 disponibili. Limite dichiarato "
        "esplicitamente.",
        "Il quarto metodo (Diffusion + DPS) e' stato scartato per onere computazionale "
        "-- opzione esplicitamente consentita dal testo per un progetto da uno studente.",
    ])

    # === Slide 3: Tre Approcci -- fix imprecise "no physics at all" claim ===
    # shape 115 has SHAPE_TO_FIT_TEXT and nothing below it in its column: safe to grow.
    s3 = slides[2]
    set_run_text(s3, 115, 0, 0,
        "La rete apprende l'intera mappa diretta dai dati degradati alle immagini "
        "pulite, senza mai usare esplicitamente l'operatore A. Eccezione: "
        "l'apprendimento residuale (x̂=y+UNet(y)) assume che y sia gia' vicina a x -- "
        "una forma debole di prior.")

    # === Slide 8: FISTA implementation -- add convergence check to last bullet ===
    # shape 176 (Discrepanza) is last, bottom=5.0in, 2.5in free below.
    s8 = slides[7]
    append_run_text(s8, 176, 0, 1,
        " Verifica di convergenza (fatta a posteriori): il PSNR vs verita' a terra "
        "picca gia' a ~10-25 iterazioni e scende lievemente fino a un plateau verso "
        "150-200 -- 'semiconvergenza', tipica dei metodi iterativi di "
        "regolarizzazione. Le 100 iterazioni scelte servono alla convergenza "
        "dell'ottimizzazione, non a massimizzare il PSNR.")

    # === Slide 9: PD-Net -- fix epoch count 20 -> 12 (exact swap), add batch size ===
    s9 = slides[8]
    set_run_text(s9, 190, 0, 2, "12")
    append_run_text(s9, 190, 0, 3, " Batch size 8.")

    # === Slide 10: UNet -- add batch size to last bullet (epoch count already correct) ===
    s10 = slides[9]
    append_run_text(s10, 205, 0, 1, " Batch size 8.")

    # === Slide 11: table -- mean+-std instead of mean only; note on SSIM below table ===
    s11 = slides[10]
    table = get_shape(s11, 217).table
    stats = {
        ("FISTA-Wavelet", "0.005"): (31.17, 3.85, 0.861, 0.068),
        ("FISTA-Wavelet", "0.01"): (31.00, 3.87, 0.855, 0.071),
        ("FISTA-Wavelet", "0.05"): (28.96, 3.35, 0.779, 0.082),
        ("FISTA-Wavelet", "0.1"): (27.09, 2.97, 0.699, 0.092),
        ("PD-Net", "0.005"): (33.43, 3.51, 0.904, 0.050),
        ("PD-Net", "0.01"): (33.84, 3.64, 0.910, 0.053),
        ("PD-Net", "0.05"): (31.77, 3.43, 0.862, 0.073),
        ("PD-Net", "0.1"): (30.14, 3.27, 0.822, 0.088),
        ("UNet", "0.005"): (33.39, 3.51, 0.903, 0.049),
        ("UNet", "0.01"): (34.03, 3.73, 0.913, 0.052),
        ("UNet", "0.05"): (32.08, 3.44, 0.872, 0.069),
        ("UNet", "0.1"): (30.49, 3.26, 0.836, 0.080),
    }
    table.cell(0, 2).text = "PSNR (dB), media±std"
    table.cell(0, 3).text = "SSIM, media±std"
    for r in range(1, 13):
        method, level = table.cell(r, 0).text, table.cell(r, 1).text
        psnr_m, psnr_s, ssim_m, ssim_s = stats[(method, level)]
        table.cell(r, 2).text = f"{psnr_m:.2f}±{psnr_s:.2f}"
        table.cell(r, 3).text = f"{ssim_m:.3f}±{ssim_s:.3f}"
    for c in (2, 3):
        for r in range(13):
            run = table.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = "Calibri"
            run.font.size = Pt(10.5 if r == 0 else 10)
            run.font.bold = (r == 0)
    add_textbox(s11, 0.75, 5.3, 11.6, 0.6, [
        "SSIM calcolato su RGB con skimage (channel_axis=2): media dell'indice SSIM "
        "sui 3 canali colore, calcolati indipendentemente.",
    ], size=12)

    # === Slide 14: Anomalia -- add FISTA-monotonic point, tight budget (~1.6in free) ===
    s14 = slides[13]
    shape238 = get_shape(s14, 238)
    p_new = shape238.text_frame.add_paragraph()
    r_new = p_new.add_run()
    r_new.text = (
        "FISTA, senza training, e' invece perfettamente monotona (31.17 -> 31.00 -> "
        "28.96 -> 27.09 dB): l'anomalia e' confinata ai metodi allenati, il che "
        "rafforza l'ipotesi che sia un effetto del training. Nota: le curve di "
        "validazione per epoca non sono state salvate su Colab, quindi non possiamo "
        "confermare se il modello a sigma=0.005 avesse gia' raggiunto un plateau."
    )
    style_run(r_new, size=13)

    # === Slide 15: Efficienza PD-Net -- fix param count (short, tight gap to shape 246) ===
    s15 = slides[14]
    set_run_text(s15, 245, 0, 0,
        "PD-Net eguaglia l'accuratezza di UNet con ~120 volte meno parametri "
        "(69.700 contro 8.321.619 -- conteggio reale dei pesi, non la dimensione del "
        "file). La conoscenza fisica sostituisce enormi quantita' di dati e pesi.")
    # shape 246 is last on this slide, 2.27in free below: safe for the longer, softened text.
    set_run_text(s15, 246, 0, 0,
        "In letteratura, un inductive bias cosi' forte e' associato a una migliore "
        "generalizzazione fuori distribuzione -- claim della letteratura, non "
        "verificata qui (nessun test out-of-distribution condotto). Cio' che e' "
        "verificato: a sigma=0.005 PD-Net (33.43 dB) e UNet (33.39 dB) differiscono "
        "di soli 0.04 dB, ben dentro la deviazione standard (~3.5 dB su 80 immagini) "
        "-- i due metodi sono staticamente indistinguibili in accuratezza.")

    # === Slide 16: Conclusioni -- fix ratio (short), hardware note on last bullet ===
    s16 = slides[15]
    set_run_text(s16, 254, 0, 1,
        " PD-Net dimostra che incorporare la fisica nota riduce drasticamente il "
        "numero di parametri necessari (~120x: 69.700 contro 8.321.619 di UNet).")
    # shape 256 is last, 2.5in free below: safe to append the hardware caveat there.
    append_run_text(s16, 256, 0, 4,
        " Nota metodologica: FISTA e' stato eseguito su CPU, UNet e PD-Net in "
        "inferenza su GPU (MPS) -- il confronto di velocita' e' quindi indicativo, "
        "non a parita' di hardware.")

    # === New slide: perche' 4 modelli specializzati ===
    template = slides[13]  # anomaly slide: title + single autofit body, dark bg
    s_spec = duplicate_slide(prs, template)
    set_paragraph_text(get_shape(s_spec, 237), 0, "Perche' 4 modelli specializzati",
        size=TITLE_SIZE, bold=True, color=TITLE_COLOR, font=TITLE_FONT)
    body = get_shape(s_spec, 238)
    tf = body.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    set_paragraph_text(body, 0,
        "FISTA specializza lambda per livello di rumore per necessita' matematica "
        "(principio di discrepanza): per un confronto equo, diamo lo stesso 'budget "
        "di specializzazione' anche a UNet e PD-Net, allenando 4 modelli indipendenti "
        "invece di uno generalista.", size=BODY_SIZE, color=BODY_COLOR, font=BODY_FONT)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = (
        "Assunzione dichiarata: il livello di rumore deve essere noto a tempo di "
        "test (si sceglie il checkpoint sapendo sigma in anticipo) -- un'informazione "
        "laterale che un metodo 'blind' non avrebbe. Un modello generalista sarebbe "
        "piu' realistico ma verosimilmente meno accurato a ciascun livello: "
        "compromesso lasciato a lavoro futuro."
    )
    style_run(r2)

    # === New slide: assunzioni e limiti ===
    s_lim = duplicate_slide(prs, template)
    set_paragraph_text(get_shape(s_lim, 237), 0, "Assunzioni e limiti dichiarati",
        size=TITLE_SIZE, bold=True, color=TITLE_COLOR, font=TITLE_FONT)
    body2 = get_shape(s_lim, 238)
    tf2 = body2.text_frame
    for p in list(tf2.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    set_paragraph_text(body2, 0,
        "Inverse crime (lieve): lo stesso identico operatore A genera i dati "
        "degradati ed e' usato dentro FISTA/PD-Net per ricostruire -- previsto e "
        "discusso nel corso, ma va dichiarato esplicitamente.",
        size=BODY_SIZE, color=BODY_COLOR, font=BODY_FONT)
    for txt in [
        "Sfocatura intrinseca del dataset: le immagini KaoKore sono ritagli di "
        "dipinti storici ridimensionati a 256x256 -- la 'ground truth' non e' "
        "perfettamente nitida in partenza, limite del dataset piu' che del metodo.",
        "Hardware non omogeneo nei tempi riportati: FISTA su CPU, UNet/PD-Net in "
        "inferenza su GPU (MPS) -- il confronto '40x piu' lento' e' indicativo.",
        "Curve di validazione per epoca non salvate durante il training su Colab: "
        "non verificabile a posteriori se un modello avesse raggiunto un plateau.",
    ]:
        p = tf2.add_paragraph()
        r = p.add_run()
        r.text = txt
        style_run(r)

    # === New slide: crop + difference map ===
    s_crop = duplicate_slide(prs, slides[12])  # visual comparison slide as layout template
    set_run_text(s_crop, 230, 0, 0, "Crop ingrandito e mappe di errore")
    get_shape(s_crop, 231)._element.getparent().remove(get_shape(s_crop, 231)._element)
    s_crop.shapes.add_picture("results/figures/crop_diff_comparison.png",
        Inches(1.0), Inches(1.5), width=Inches(11.3))

    # === New slide: bibliography ===
    s_bib = duplicate_slide(prs, template)
    set_paragraph_text(get_shape(s_bib, 237), 0, "Bibliografia",
        size=TITLE_SIZE, bold=True, color=TITLE_COLOR, font=TITLE_FONT)
    body3 = get_shape(s_bib, 238)
    tf3 = body3.text_frame
    for p in list(tf3.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    refs = [
        "A. Beck, M. Teboulle. A Fast Iterative Shrinkage-Thresholding Algorithm for "
        "Linear Inverse Problems. SIAM J. Imaging Sciences, 2009.",
        "A. Chambolle, T. Pock. A First-Order Primal-Dual Algorithm for Convex "
        "Problems with Applications to Imaging. JMIV, 2011.",
        "J. Adler, O. Oktem. Learned Primal-Dual Reconstruction. IEEE Trans. Medical "
        "Imaging, 2018.",
        "O. Ronneberger, P. Fischer, T. Brox. U-Net: Convolutional Networks for "
        "Biomedical Image Segmentation. MICCAI, 2015.",
        "Y. Tian et al. KaoKore: A Pre-modern Japanese Art Facial Expression Dataset "
        "(dataset), 2020.",
        "D. Evangelista. IPPy library (materiale del corso), da cui sono riadattati "
        "gli operatori Blurring/Gradient e lo schema Learned Primal-Dual.",
    ]
    set_paragraph_text(body3, 0, refs[0], size=13, color=BODY_COLOR, font=BODY_FONT)
    for txt in refs[1:]:
        p = tf3.add_paragraph()
        r = p.add_run()
        r.text = txt
        style_run(r, size=13)

    prs.save(PATH)
    print("saved", PATH, "- slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
