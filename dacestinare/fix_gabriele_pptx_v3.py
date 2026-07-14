"""Third pass: fixes three text overlaps introduced when boxes grew past their fixed
siblings. (1) Table slide: the baseline note collided with the taller table -> remove
it (baseline is shown as a curve in the plot). (2) Efficiency slide: two boxes
overlapped -> merge into one autofitting box. (3) Conclusions: first bullet grew into
the second -> shorten it (the efficiency detail is already the second bullet)."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

PATH = "presentation/gabriele_centonze.pptx"
BODY_FONT, BODY_COLOR = "Open Sans", RGBColor(0xCB, 0xD5, 0xE1)


def title_of(slide):
    best = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.top if sh.top is not None else 0
            if best is None or t < best[0]:
                best = (t, sh.text_frame.text)
    return best[1].replace("\x0b", " ").replace("\n", " ").strip()


def slide_by_title(prs, prefix):
    for slide in prs.slides:
        if title_of(slide).startswith(prefix):
            return slide
    raise ValueError(prefix)


def get_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(shape_id)


def main():
    prs = Presentation(PATH)

    # (1) Table slide: remove the baseline/SSIM note that overlaps the table.
    s = slide_by_title(prs, "Risultati Numerici")
    for sh in list(s.shapes):
        if sh.has_text_frame and "Baseline (osservazione" in sh.text_frame.text:
            sh._element.getparent().remove(sh._element)

    # (2) Efficiency slide: merge shape 246 into 245, then delete 246.
    s = slide_by_title(prs, "L'Efficienza di PD-Net")
    box245 = get_shape(s, 245)
    box246 = get_shape(s, 246)
    text246 = box246.text_frame.paragraphs[0].runs[0].text
    p = box245.text_frame.add_paragraph()
    r = p.add_run()
    r.text = text246
    r.font.size = Pt(15)
    r.font.color.rgb = BODY_COLOR
    r.font.name = BODY_FONT
    box246._element.getparent().remove(box246._element)

    # (3) Conclusions: shorten the first bullet so it no longer grows into the second.
    s = slide_by_title(prs, "Conclusioni")
    get_shape(s, 253).text_frame.paragraphs[0].runs[1].text = (
        " UNet e PD-Net battono FISTA di 2-3 dB a ogni livello; il prior appreso su "
        "KaoKore supera la sparsita' wavelet generica. Tra i due, UNet e' leggermente "
        "avanti.")

    prs.save(PATH)
    print("saved", PATH)


if __name__ == "__main__":
    main()
