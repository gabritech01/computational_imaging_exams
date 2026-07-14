"""Second round of corrections on gabriele_centonze.pptx, applying an external
technical review. Starts from the current 20-slide file (backed up separately).

Fixes factual errors (semiconvergence misused, TV proximals wrongly called
'without closed form', wrong paired-statistics claim, over-promised 'equita'),
rephrases the non-monotonic anomaly honestly, adds the degraded baseline, adds a
'Setup sperimentale' slide required by the track, and reorders so Conclusions is
last (before Bibliography) and the error maps sit with the other results.

Language kept deliberately plain (exam support), without over-simplifying the
concepts. Slides are addressed by title so edits survive the reordering.
"""
import copy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

PATH = "presentation/gabriele_centonze.pptx"

TITLE_FONT, TITLE_SIZE, TITLE_COLOR = "Merriweather", 24, RGBColor(0x10, 0xB9, 0x81)
BODY_FONT, BODY_SIZE, BODY_COLOR = "Open Sans", 15, RGBColor(0xCB, 0xD5, 0xE1)

TARGET_ORDER = [
    "Deblur & Denoise su KaoKore",
    "Descrizione del Progetto",
    "Tre Approcci al Problema Inverso",
    "Il Filo Conduttore",
    "Setup sperimentale",
    "Metodologia: Variazionale (FISTA)",
    "Metodologia: Ibrido (PD-Net)",
    "Metodologia: End-to-end (UNet)",
    "Implementazione: FISTA + Wavelet",
    "Implementazione: PD-Net",
    "Implementazione: UNet",
    "Perche' 4 modelli specializzati",
    "Risultati Numerici: Tabella Riassuntiva",
    "Confronto Prestazionale",
    "Confronto Visivo",
    "Crop ingrandito e mappe di errore",
    "L'Anomalia Non Monotona",
    "L'Efficienza di PD-Net",
    "Assunzioni e limiti dichiarati",
    "Conclusioni e Sviluppi Futuri",
    "Bibliografia",
]


def get_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"shape {shape_id} not found")


def set_run_text(slide, shape_id, para_idx, run_idx, text):
    get_shape(slide, shape_id).text_frame.paragraphs[para_idx].runs[run_idx].text = text


def style_run(run, size=BODY_SIZE, color=BODY_COLOR, font=BODY_FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font


def title_of(slide):
    best = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.top if sh.top is not None else 0
            if best is None or t < best[0]:
                best = (t, sh.text_frame.text)
    return best[1].replace("\x0b", " ").replace("\n", " ").strip()


def slide_by_title(prs, prefix):
    for slide in prs.slides:
        if title_of(slide).startswith(prefix):
            return slide
    raise ValueError(f"slide not found: {prefix}")


def find_shape_by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    raise ValueError(f"shape with text {needle!r} not found")


def duplicate_slide(prs, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    rid_map = {}
    for rId, rel in source_slide.part.rels.items():
        if rel.reltype.endswith("slideLayout") or rel.reltype.endswith("notesSlide"):
            continue
        if rel.is_external:
            rid_map[rId] = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            rid_map[rId] = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        for el in new_el.iter():
            for attr in ("embed", "link"):
                val = el.get(qn("r:" + attr))
                if val in rid_map:
                    el.set(qn("r:" + attr), rid_map[val])
        new_slide.shapes._spTree.append(new_el)
    return new_slide


def set_paragraph_text(shape, para_idx, text, size=None, bold=None, color=None, font=None):
    p = shape.text_frame.paragraphs[para_idx]
    if not p.runs:
        p.add_run()
    r = p.runs[0]
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    r.text = text
    if size:
        r.font.size = Pt(size)
    if bold is not None:
        r.font.bold = bold
    if color:
        r.font.color.rgb = color
    if font:
        r.font.name = font


def reorder(prs):
    sldIdLst = prs.slides._sldIdLst
    id_elems = list(sldIdLst)
    pairs = [(title_of(slide), ie) for slide, ie in zip(list(prs.slides), id_elems)]
    used, ordered = set(), []
    for target in TARGET_ORDER:
        for i, (title, ie) in enumerate(pairs):
            if i not in used and title.startswith(target):
                ordered.append(ie)
                used.add(i)
                break
        else:
            raise ValueError(f"reorder: title not found {target!r}")
    assert len(ordered) == len(id_elems), "reorder count mismatch"
    for ie in id_elems:
        sldIdLst.remove(ie)
    for ie in ordered:
        sldIdLst.append(ie)


def main():
    prs = Presentation(PATH)

    # --- Slide Descrizione: soften over-promised 'equita' + augment bottom note ---
    s = slide_by_title(prs, "Descrizione del Progetto")
    set_run_text(s, 95, 0, 1,
        " Implementare e confrontare criticamente metodi di famiglie diverse, con la "
        "stessa degradazione e lo stesso test set per tutti (confronto equo sui dati).")
    note = find_shape_by_text(s, "Nota sui numeri")
    p = note.text_frame.add_paragraph()
    r = p.add_run()
    r.text = ("I budget di training dei metodi appresi differiscono (PD-Net 700 img/12 "
              "epoche, UNet 1000/20): la parita' e' quindi raggiunta da PD-Net con meno "
              "risorse, il che rafforza la tesi sull'efficienza.")
    style_run(r, size=13)

    # --- Slide Tre Approcci: fix 'proximali senza forma chiusa' on the Ibrido card ---
    s = slide_by_title(prs, "Tre Approcci")
    set_run_text(s, 113, 0, 0,
        "Combina la fisica esatta (operatore A e gradiente) con piccole reti neurali "
        "che imparano dai dati un regolarizzatore piu' ricco della sola TV, mantenendo "
        "esatta la parte fisica.")

    # --- Slide Metodologia PD-Net: same fix, and nudge the bullet below to keep spacing ---
    s = slide_by_title(prs, "Metodologia: Ibrido")
    set_run_text(s, 149, 0, 0,
        "I proximali della TV hanno gia' forma chiusa nota: li sostituiamo con piccole "
        "CNN per imparare un regolarizzatore piu' espressivo della TV.")
    b150 = get_shape(s, 150)
    b150.top = b150.top + Inches(0.35)

    # --- Slide Impl FISTA: replace the misused 'semiconvergenza', add step size + prox note ---
    s = slide_by_title(prs, "Implementazione: FISTA")
    set_run_text(s, 176, 0, 1,
        " I valori di lambda ottimali (0.0025, 0.0035, 0.0175, 0.07) crescono col "
        "rumore, in accordo col principio di discrepanza. Passo di FISTA: 1/L, con "
        "L=||A^T A||=1 grazie al kernel normalizzato. Con W ortogonale il prossimale "
        "del termine L1 e' esattamente il soft-thresholding dei coefficienti wavelet. "
        "Il PSNR non e' monotono lungo le iterazioni: picca prima e poi si assesta sul "
        "valore del minimizzatore del funzionale -- non e' semiconvergenza (il problema "
        "regolarizzato e' ben posto e l'errore non diverge), ma la differenza tra "
        "minimizzare il funzionale e massimizzare una metrica esterna come il PSNR.")

    # --- Slide Tabella: add degraded baseline to the note below the table ---
    s = slide_by_title(prs, "Risultati Numerici")
    ssim_note = find_shape_by_text(s, "SSIM calcolato")
    first_run = ssim_note.text_frame.paragraphs[0].runs[0]
    first_run.text = (
        "Baseline (osservazione degradata, PSNR): 29.9 / 29.8 / 26.7 / 21.6 dB ai 4 "
        "livelli -- il guadagno di ogni metodo si misura da qui (vedi anche il plot). "
        + first_run.text)

    # --- Slide Anomalia: honest hypotheses instead of the weak 'tiny correction' story ---
    s = slide_by_title(prs, "L'Anomalia Non Monotona")
    body = get_shape(s, 238)
    tf = body.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    paras = [
        "Sia UNet che PD-Net raggiungono un PSNR leggermente piu' basso a sigma=0.005 "
        "che a sigma=0.01: controintuitivo, il caso meno rumoroso va peggio.",
        "L'input degradato e' pero' quasi identico ai due livelli (baseline 29.86 vs "
        "29.81 dB): il blur domina sul rumore, quindi la difficolta' del compito cambia "
        "pochissimo. La causa non sta nei dati.",
        "FISTA (senza training) e' invece perfettamente monotona: l'anomalia riguarda "
        "solo i metodi allenati, quindi e' un effetto del training, non del problema.",
        "Ipotesi, dalla piu' probabile: (i) la selezione del modello e' fatta sul dev "
        "set, ma piccolo (20 immagini), quindi rumorosa; (ii) un solo run di training "
        "per configurazione (nessuna media su piu' seed); (iii) la loss L1 non e' "
        "perfettamente allineata col PSNR (che e' L2). Servirebbero piu' run e un dev "
        "set piu' grande per distinguerle.",
    ]
    set_paragraph_text(body, 0, paras[0], size=14, color=BODY_COLOR, font=BODY_FONT)
    for txt in paras[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        style_run(r, size=14)

    # --- Slide Efficienza PD-Net: fix 'eguaglia' + wrong statistics; drop OOD non-sequitur ---
    s = slide_by_title(prs, "L'Efficienza di PD-Net")
    set_run_text(s, 245, 0, 0,
        "PD-Net resta entro ~0.4 dB da UNet usando circa 120 volte meno parametri "
        "(69.700 contro 8.321.619 -- conteggio reale dei pesi, non la dimensione del "
        "file). Incorporare la fisica nota (A e gradiente) riduce moltissimo cio' che "
        "la rete deve imparare.")
    set_run_text(s, 246, 0, 0,
        "Confronto appaiato sulle stesse 80 immagini: UNet e' significativamente "
        "migliore ai 3 livelli di rumore piu' alti, ma con distacchi piccoli "
        "(0.16-0.39 dB); a sigma=0.005 i due sono indistinguibili (-0.04 dB). Il punto "
        "non e' chi vince, ma che PD-Net arriva a una frazione di dB da una rete 120x "
        "piu' grande e allenata con piu' dati.")

    # --- Slide Conclusioni: keep headline honest (UNet leggermente avanti) ---
    s = slide_by_title(prs, "Conclusioni")
    set_run_text(s, 253, 0, 1,
        " UNet e PD-Net battono FISTA di 2-3 dB a ogni livello. Il prior appreso sulla "
        "statistica di KaoKore supera la sparsita' wavelet generica. UNet e' "
        "leggermente avanti, PD-Net resta entro 0.4 dB con 120x meno parametri.")

    # --- New slide: Setup sperimentale (required by the track) ---
    template = slide_by_title(prs, "L'Anomalia Non Monotona")
    s_setup = duplicate_slide(prs, template)
    set_paragraph_text(get_shape(s_setup, 237), 0, "Setup sperimentale",
        size=TITLE_SIZE, bold=True, color=TITLE_COLOR, font=TITLE_FONT)
    body = get_shape(s_setup, 238)
    body.top = Inches(1.5)
    body.height = Inches(5.2)
    tf = body.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    setup_paras = [
        "Immagini normalizzate in [0,1], RGB (float); il blur e' applicato in modo "
        "indipendente sui 3 canali colore.",
        "Blur: kernel gaussiano 9x9, sigma=2, condizioni al bordo circolari (periodiche) "
        "-- scelta che rende A e il suo aggiunto A^T esatti e diagonalizzabili via FFT.",
        "Rumore: gaussiano additivo i.i.d., con deviazione standard pari al livello di "
        "rumore sulla scala [0,1], per-pixel e per-canale. L'osservazione viene salvata "
        "come PNG a 8 bit (quantizzazione, come un sensore reale) e riusata identica da "
        "tutti i metodi.",
        "Metriche: PSNR e SSIM con data_range=1; SSIM su RGB come media dei 3 canali.",
        "Split effettivo usato: 700-1000 train / 20 dev / 80 test, sottoinsieme fisso "
        "dello split ufficiale KaoKore (scelto per vincoli di tempo).",
    ]
    set_paragraph_text(body, 0, setup_paras[0], size=15, color=BODY_COLOR, font=BODY_FONT)
    for txt in setup_paras[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        style_run(r, size=15)

    reorder(prs)

    prs.save(PATH)
    print("saved", PATH, "- slides:", len(prs.slides._sldIdLst))
    print("order:")
    for i, slide in enumerate(prs.slides):
        print(f"  {i+1}. {title_of(slide)[:45]}")


if __name__ == "__main__":
    main()
