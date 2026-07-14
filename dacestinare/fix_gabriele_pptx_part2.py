"""Second pass of corrections on gabriele_centonze.pptx: adds the FISTA convergence
check finding (semiconvergence, PSNR peaks well before 100 iterations) and clarifies
how SSIM is computed on RGB images -- both flagged as missing technical details.
Run only once, after fix_gabriele_pptx.py (not idempotent, appends text)."""
import shutil

from pptx import Presentation
from pptx.util import Pt

PATH = "presentation/gabriele_centonze.pptx"
BACKUP = "presentation/old/gabriele_centonze_before_fix2.pptx"


def get_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"shape {shape_id} not found")


def append_run_text(slide, shape_id, para_idx, run_idx, extra):
    shape = get_shape(slide, shape_id)
    r = shape.text_frame.paragraphs[para_idx].runs[run_idx]
    r.text = r.text + extra


def main():
    shutil.copy(PATH, BACKUP)
    prs = Presentation(PATH)
    slides = list(prs.slides)

    # Slide 8 (FISTA implementation): convergence check on the "Ciclo" bullet
    s8 = slides[7]
    append_run_text(s8, 174, 0, 1,
        " Verifica di convergenza fatta a posteriori su un'immagine test: il PSNR "
        "rispetto alla verita' a terra in realta' picca gia' a ~10-25 iterazioni "
        "(33.5 dB) e scende leggermente fino a un plateau intorno a 150-200 "
        "iterazioni (31.7 dB) -- fenomeno di 'semiconvergenza' tipico dei metodi "
        "iterativi di regolarizzazione: il minimizzatore esatto del funzionale "
        "penalizzato non coincide con la verita' a terra. 100 iterazioni sono state "
        "scelte per la convergenza dell'ottimizzazione, non perche' massimizzino il "
        "PSNR -- un margine di miglioramento lasciato a lavoro futuro.")

    # Slide 11 (results table): note on how SSIM is computed on RGB
    s11 = slides[10]
    title_shape = get_shape(s11, 216)
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ("SSIM calcolato su RGB con skimage (channel_axis=2): media dell'indice "
              "SSIM calcolato indipendentemente su ciascuno dei 3 canali colore.")
    r.font.size = Pt(13)
    r.font.italic = True

    prs.save(PATH)
    print("saved", PATH)


if __name__ == "__main__":
    main()
